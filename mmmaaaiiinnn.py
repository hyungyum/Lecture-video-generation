# -*- coding: utf-8 -*-
import os
import sys
import locale
import uuid
import tempfile
import subprocess
import io
import base64
import json
import requests

import streamlit as st
from streamlit_option_menu import option_menu
from PyPDF2 import PdfReader
from pptx import Presentation
import fitz  # PyMuPDF for PDF text extraction
import openai
from openai import OpenAI
from pdf2image import convert_from_path, convert_from_bytes
from PIL import Image
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips

import torch
import torchaudio
from zonos.model import Zonos, DEFAULT_BACKBONE_CLS as ZonosBackbone
from zonos.conditioning import make_cond_dict, supported_language_codes
from zonos.utils import DEFAULT_DEVICE as device

from pydantic import BaseModel
from typing import List
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv
from collections import defaultdict
from pprint import pprint

# ─────────────── 초기 설정 ───────────────
load_dotenv()  # .env 파일에 OPENAI_API_KEY=sk-... 설정

# UTF-8 강제 인코딩
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
    sys.stdin.reconfigure(encoding="utf-8")
os.environ["PYTHONUTF8"] = "1"
os.environ["PYTHONIOENCODING"] = "utf-8"
locale.setlocale(locale.LC_ALL, 'C.UTF-8')

# OpenAI 키 설정
openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key or not openai.api_key.isascii():
    raise ValueError("유효한 OPENAI_API_KEY가 필요합니다.")

# GPT-4o 클라이언트 생성
client = OpenAI(api_key=openai.api_key)




class SlideModel(BaseModel):
    title: str
    bullets: List[str]
    narration: str


# ─────────────── 유틸 함수 ───────────────

def extract_from_pptx(data: bytes):
    prs = Presentation(io.BytesIO(data))
    slides = []
    for slide in prs.slides:
        # 텍스트 추출
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        text = "\n".join(text_parts).strip()
        # 이미지 추출
        imgs = []
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                # shape_type==13은 Picture
                try:
                    blob = shape.image.blob
                    imgs.append(blob)
                except Exception:
                    continue
        slides.append({"text": text, "images": imgs})
    return slides


def extract_from_pdf(data: bytes):
    doc = fitz.open(stream=data, filetype="pdf")
    slides = []
    for page in doc:
        text = page.get_text("text", sort=True).strip()
        # 페이지 전체를 하나의 이미지로 변환
        pix = page.get_pixmap()
        img_bytes = pix.tobytes()
        slides.append({"text": text, "images": [img_bytes]})
    doc.close()
    return slides


def extract_slide_texts(ppt_path: str) -> list[str]:
    prs = Presentation(ppt_path)
    slides = []
    for slide in prs.slides:
        parts = []
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False) and sh.text.strip():
                parts.append(sh.text.strip())
        note = ""
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
        slides.append(note if note else "\n".join(parts))
    return slides


def extract_pdf_texts(pdf_path: str) -> list[str]:
    doc = fitz.open(pdf_path)
    texts = []
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        txt = page.get_text("text", sort=True).strip()
        texts.append(txt)
    doc.close()
    return texts





def convert_slides_to_images_auto(input_fp: str, out_dir: str, dpi: int = 200) -> list[str]:
    ext = os.path.splitext(input_fp)[1].lower()
    os.makedirs(out_dir, exist_ok=True)
    saved = []

    if ext == ".pptx":
        try:
            st.info("PPTX를 PDF로 변환 중...")
            pdf_dir = os.path.dirname(input_fp)
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", pdf_dir, input_fp],
                check=True
            )
            pdf_path = os.path.join(pdf_dir, os.path.basename(input_fp).rsplit(".pptx",1)[0] + ".pdf")
            
            st.info("PDF를 이미지로 변환 중...")
            images = convert_from_path(pdf_path, dpi=dpi, fmt="png", paths_only=False)
            for i, img in enumerate(images, start=1):
                fn = os.path.join(out_dir, f"slide_{i:03}.png")
                img.save(fn, "PNG")
                saved.append(fn)
            st.success(f"{len(saved)}개 슬라이드를 이미지로 변환했습니다.")
            
            # 임시 PDF 파일 삭제
            os.remove(pdf_path)
        except Exception as e:
            st.error(f"PPTX 변환 오류: {e}")
            # 매우 간단한 이미지 생성 (최후의 수단)
            if 'texts' in locals() and texts:
                slide_count = len(texts)
            else:
                from pptx import Presentation
                try:
                    pres = Presentation(input_fp)
                    slide_count = len(pres.slides)
                except:
                    slide_count = 5  # 기본값
                    
            for i in range(1, slide_count+1):
                img = Image.new('RGB', (1280, 720), (255, 255, 255))
                draw = ImageDraw.Draw(img)
                draw.text((640, 360), f"Slide {i}", fill=(0, 0, 0), anchor="mm")
                fn = os.path.join(out_dir, f"slide_{i:03}.png")
                img.save(fn, "PNG")
                saved.append(fn)
    else:
        # PDF → 이미지
        try:
            st.info("PDF를 이미지로 변환 중...")
            # pdf2image 사용 (poppler-utils 필요, packages.txt에 추가됨)
            images = convert_from_path(input_fp, dpi=dpi, fmt="png")
            for i, img in enumerate(images, start=1):
                fn = os.path.join(out_dir, f"slide_{i:03}.png")
                img.save(fn, "PNG")
                saved.append(fn)
            st.success(f"{len(saved)}개 PDF 페이지를 변환했습니다.")
        except Exception as e:
            st.warning(f"pdf2image 실패: {e}")
            try:
                # PyMuPDF 사용
                st.info("PyMuPDF로 PDF 변환 중...")
                import fitz
                doc = fitz.open(input_fp)
                for i, page in enumerate(doc, start=1):
                    pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
                    fn = os.path.join(out_dir, f"slide_{i:03}.png")
                    pix.save(fn)
                    saved.append(fn)
                doc.close()
                st.success(f"{len(saved)}개 PDF 페이지를 변환했습니다.")
            except Exception as fitz_err:
                st.error(f"PDF 변환 실패: {fitz_err}")

    return saved


def extract_text_from_pdf(pdf_reader, start_page, end_page):
    text = []
    for i in range(start_page - 1, end_page):
        page = pdf_reader.pages[i]
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)
    return "\n".join(text)


def extract_chapter_pages(reader):
    """PDF에서 챕터별 시작 페이지와 끝 페이지 추출"""
    total_pages = len(reader.pages)
    outlines = reader.outline

    if not outlines:
        return None, []

    # 1. 평면화된 목차 구하기
    flat_outlines = []

    def flatten_with_level(items, level=0):
        for item in items:
            if isinstance(item, list):
                flatten_with_level(item, level + 1)
            else:
                try:
                    page_num = reader.get_destination_page_number(item)
                    flat_outlines.append(
                        {
                            "title": item.title,
                            "page": page_num + 1,  # 0-based에서 1-based로 변환
                            "level": level,
                        }
                    )
                except Exception:
                    pass

    flatten_with_level(outlines)
    flat_outlines.sort(key=lambda x: x["page"])  # 페이지 순으로 정렬

    # 2. 챕터와 서브챕터 구조 만들기
    chapters = []
    for item in flat_outlines:
        if item["level"] == 0:  # 챕터 레벨
            chapters.append(
                {"title": item["title"], "page": item["page"], "subchapters": []}
            )
        elif item["level"] == 1 and chapters:  # 서브챕터이고 최소 하나의 챕터가 있는 경우
            chapters[-1]["subchapters"].append(
                {"title": item["title"], "page": item["page"]}
            )

    return chapters, flat_outlines


# ─────────────── 사이드바에 옵션 메뉴 추가 ───────────────────────────────────
with st.sidebar:
    selected = option_menu(
        menu_title="네비게이션",
        options=["강의 영상 생성", "From PDF Book", "From Slide"], # 강의 영상 생성 , 강의 자료 생성
        icons=[
            "house",
            "file-ppt",
            "camera-video",
        ],
        menu_icon="app-indicator",
        default_index=0,
        orientation="vertical",
        styles={
            "container": {
                "padding": "5!important",
                "background-color": "lightgray",
            },
            "icon": {"color": "black", "font-size": "18px"},
            "nav-link": {
                "color": "black",
                "font-size": "20px",
                "text-align": "left",
                "margin": "0px",
                "--hover-color": "#f0f0f0",
            },
            "nav-link-selected": {"background-color": "#02ab21"},
        },
    )



if selected == "강의 영상 생성":

    st.markdown(
    """
    <style>
    /* 페이지 패딩 제거 */
    section.main .block-container { padding: 0; margin: 0; }
    /* 이미지가 화면 꽉 차게 */
    .stImage > img {
        width: 100vw !important;
        height: 100vh !important;
        object-fit: cover;
    }
    </style>
    """,
    unsafe_allow_html=True,
)
    st.image("system6.png")


    # 샘플 파일 목록 정의
    sample_ppts = ["딥러닝.pptx", "미국의 역사.pptx"]
    sample_videos = ["딥러닝.mp4", "미국의 역사.mp4"]

    st.markdown("**생성된 예시**")
    col1, col2 = st.columns(2)

    # 1) 세션에 기본값 설정 (최초 실행 시)
    if "ppt_choice" not in st.session_state:
        st.session_state.ppt_choice = sample_ppts[0]

    if "vid_choice" not in st.session_state:
        st.session_state.vid_choice = sample_videos[0]

    with col1:
        # selectbox 에 key 지정 → 세션 상태로 자동 저장
        st.selectbox("샘플 PPT 선택", sample_ppts, key="ppt_choice")
        ppt_path = os.path.join("./samples", st.session_state.ppt_choice)
        with open(ppt_path, "rb") as f:
            ppt_bytes = f.read()
        st.download_button(
            label="PPT 다운로드",
            data=ppt_bytes,
            file_name=st.session_state.ppt_choice,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    with col2:
        st.selectbox("샘플 영상 선택", sample_videos, key="vid_choice")
        vid_path = os.path.join("./samples", st.session_state.vid_choice)
        st.video(vid_path)



# ─── 선택된 메뉴에 따라 화면 분기 ─────────────────────────────────
elif selected == "From PDF Book": # 강의 자료 생성
    uploaded_file = st.file_uploader("PDF 파일 업로드", type=["pdf"])

    if uploaded_file:
        try:
            reader = PdfReader(uploaded_file)
            total_pages = len(reader.pages)

            # 목차 기반 페이지 범위 선택 시도
            try:
                chapters, flat_outlines = extract_chapter_pages(reader)

                if chapters:
                    st.info("목차를 기반으로 페이지 범위를 선택합니다.")

                    # 챕터 선택 (항상 먼저 선택)
                    chapter_options = [
                        (f"{c['title']} (p.{c['page']})", i)
                        for i, c in enumerate(chapters)
                    ]
                    selected_chapter_idx = st.selectbox(
                        "챕터 선택",
                        options=[opt[1] for opt in chapter_options],
                        format_func=lambda x: chapter_options[x][0],
                    )

                    selected_chapter = chapters[selected_chapter_idx]
                    
                    # 챕터 내 세부 범위 선택 옵션
                    selection_detail = st.radio(
                        "범위 선택",
                        ["챕터 전체 선택", "챕터 내 세부 목차 선택"],
                        horizontal=True
                    )
                    
                    if selection_detail == "챕터 전체 선택":
                        # 챕터 전체 선택 (기존 방식)
                        start_page = selected_chapter["page"]
                        
                        # 선택한 챕터의 끝 페이지 (다음 챕터 시작 전 또는 문서 끝)
                        if selected_chapter_idx < len(chapters) - 1:
                            end_page = chapters[selected_chapter_idx + 1]["page"] - 1
                        else:
                            end_page = total_pages
                            
                        st.success(
                            f"선택됨: {selected_chapter['title']} (페이지 {start_page}~{end_page})"
                        )
                        
                    else:  # 챕터 내 세부 목차 선택
                        # 현재 선택된 챕터 내의 세부 목차만 필터링
                        chapter_start_page = selected_chapter["page"]
                        
                        # 다음 챕터의 시작 페이지 또는 문서 끝
                        if selected_chapter_idx < len(chapters) - 1:
                            chapter_end_page = chapters[selected_chapter_idx + 1]["page"] - 1
                        else:
                            chapter_end_page = total_pages
                            
                        # 해당 챕터 범위 내의 세부 목차만 필터링
                        chapter_toc_options = [
                            (item["title"], item["page"]) 
                            for item in flat_outlines 
                            if chapter_start_page <= item["page"] <= chapter_end_page
                        ]
                        
                        if not chapter_toc_options:
                            st.info("이 챕터에는 세부 목차가 없습니다. 챕터 전체가 선택됩니다.")
                            start_page = chapter_start_page
                            end_page = chapter_end_page
                        else:
                            col1, col2 = st.columns(2)
                            with col1:
                                start_choice = st.selectbox(
                                    "시작 목차",
                                    options=chapter_toc_options,
                                    format_func=lambda x: f"{x[0]} (p.{x[1]})",
                                )
    
                            with col2:
                                end_choices = [
                                    opt for opt in chapter_toc_options if opt[1] >= start_choice[1]
                                ]
                                end_choice = st.selectbox(
                                    "끝 목차",
                                    options=end_choices,
                                    format_func=lambda x: f"{x[0]} (p.{x[1]})",
                                )
    
                                start_page, end_page = start_choice[1], end_choice[1]
                                
                            st.success(
                                f"선택됨: {selected_chapter['title']} 내의 {start_choice[0]} ~ {end_choice[0]} (페이지 {start_page}~{end_page})"
                            )
                else:
                    raise Exception("목차 옵션 없음")
            except Exception as e:
                st.warning("목차를 추출하지 못했습니다. 페이지 번호를 직접 입력하세요.")
                col1, col2 = st.columns(2)
                with col1:
                    start_page = st.number_input(
                        "시작 페이지", min_value=1, max_value=total_pages, value=1, step=1
                    )
                with col2:
                    end_page = st.number_input(
                        "끝 페이지",
                        min_value=start_page,
                        max_value=total_pages,
                        value=total_pages,
                        step=1,
                    )
            
            # 나레이션 난이도와 길이를 분리하여 선택
            col1, col2 = st.columns(2)
            with col1:
                difficulty_level = st.radio(
                    "나레이션 난이도",
                    options=["쉽게", "적당하게", "어렵게"],
                    index=1,
                    horizontal=True
                )
            with col2:
                length_level = st.radio(
                    "나레이션 길이",
                    options=["짧게", "적당하게", "길게"],
                    index=1,
                    horizontal=True
                )

            # 언어 옵션
            language_option = st.radio(
                "프레젠테이션 언어",
                options=[
                    "한글 슬라이드 한글 설명",
                    "영어 슬라이드 한글 설명",
                    "영어 슬라이드 영어 설명"
                ],
                index=0,
                horizontal=True
            )

            # 언어 지시사항 및 출력 언어 설정
            if language_option == "한글 슬라이드 한글 설명":
                language_instruction = "Write the narration in Korean, and use Korean slide titles and bullet points."
                output_language = "Korean"
            elif language_option == "영어 슬라이드 한글 설명":
                language_instruction = "Write the narration in Korean, but use English for slide titles and bullet points."
                output_language = "Korean"
            elif language_option == "영어 슬라이드 영어 설명":
                language_instruction = "Write the narration, slide titles, and bullet points in English."
                output_language = "English"
            else:
                language_instruction = ""
                output_language = "Korean"

            # 테마 선택 옵션 추가
            theme_option = st.radio(
                "프레젠테이션 테마",
                options=["라이트 모드", "다크 모드"],
                index=0,
                horizontal=True
            )
            
            if st.button("확인"):
                with st.spinner("텍스트를 추출하는 중..."):
                    extracted = extract_text_from_pdf(reader, start_page, end_page)
                if extracted:
                    #st.text_area("추출된 텍스트", extracted, height=400)

                    # ─── OpenAI API 호출 ───────────────────────────────
                    with st.spinner("텍스트 요약 중..."):
                        # Add difficulty instruction to system prompt
                        # 난이도 지시사항 설정
                        difficulty_instruction = ""
                        if difficulty_level == "쉽게":
                            difficulty_instruction = "Write the narration using simple words and explanations that even elementary school students can understand. Explain technical terms in simple language and use analogies and examples."
                        elif difficulty_level == "적당하게":
                            difficulty_instruction = "Write the narration at a college student level. Use technical terms appropriately, but include brief explanations when necessary."
                        elif difficulty_level == "어렵게":
                            difficulty_instruction = "Write the narration at an expert level with sophisticated vocabulary and concepts. Assume the audience has background knowledge in the subject."

                        # 길이 지시사항 설정
                        length_instruction = ""
                        if length_level == "짧게":
                            length_instruction = "Keep the narration concise and to the point. Focus only on essential information."
                        elif length_level == "적당하게":
                            length_instruction = "Use a balanced approach to length, providing sufficient detail without being excessive."
                        elif length_level == "길게":
                            length_instruction = "Provide detailed explanations with examples and elaborations on key points to ensure thorough understanding."

                        # OpenAI API 호출에 두 지시사항을 모두 포함
                        resp = client.chat.completions.create(
                            model="gpt-4.1",
                            messages=[
                                {
                                    "role": "system",
                                    "content": "You are an expert presentation designer. "
                                    "Transform the following text into an appropriate number of slides based on the content complexity. "
                                    "Each slide should have a title and concise bullet points with narration. "
                                    "Output only valid JSON matching this schema: "
                                    "[\n"
                                    '{"title": string, "bullets": [string, …], "narration": string}, '
                                    '{"title": string, "bullets": [string, …], "narration": string} ... '
                                    "]\n"
                                    "The narration should provide a sufficient explanation to help audience to understand the main idea of slide. "
                                    f"{difficulty_instruction} {length_instruction} "
                                    f"{language_instruction} "
                                    "Use your judgment to determine the optimal number of slides needed to present the content effectively. "
                                    "Do not include ''' json ''' "
                                    f"Output in {output_language}.",
                                },
                                {"role": "user", "content": extracted},
                            ],
                        )
                    ai_answer = resp.choices[0].message.content

                    #st.text_area("OpenAI 결과", ai_answer, height=300)

                    try:
                        slides_data = json.loads(ai_answer)
                        slides = [SlideModel(**item) for item in slides_data]

                        prs = Presentation()
                        # 16:9 비율
                        prs.slide_width = Inches(16)
                        prs.slide_height = Inches(9)

                        # 선택된 테마에 따른 색상 설정
                        if theme_option == "다크 모드":
                            # 다크 모드 색상
                            bg_color = RGBColor(30, 30, 38)  # 어두운 청회색
                            title_color = RGBColor(220, 220, 255)  # 밝은 청백색
                            body_color = RGBColor(200, 200, 200)  # 밝은 회색
                        else:  # 라이트 모드 (기본)
                            # 기존 라이트 모드 색상
                            bg_color = RGBColor(245, 245, 245)  # 밝은 회색
                            title_color = RGBColor(0, 51, 102)   # 짙은 파란색
                            body_color = RGBColor(60, 60, 60)    # 짙은 회색

                        for slide in slides:
                            sld = prs.slides.add_slide(prs.slide_layouts[1])
                            # 배경 색깔
                            bg = sld.background.fill
                            bg.solid()
                            bg.fore_color.rgb = bg_color

                            # 제목 설정 및 스타일
                            title_shape = sld.shapes.title
                            title_shape.text = slide.title

                            title_p = title_shape.text_frame.paragraphs[0]
                            title_p.font.name = "Arial"
                            title_p.font.size = Pt(48)
                            title_p.font.bold = True
                            title_p.font.color.rgb = title_color
                            title_p.alignment = PP_ALIGN.CENTER

                            # 본문 텍스트 프레임
                            body_shape = sld.shapes.placeholders[1]
                            tf = body_shape.text_frame
                            tf.clear()

                            cleaned = [
                                b.replace("\n", " ").strip() for b in slide.bullets
                            ]

                            if cleaned:
                                # 첫 번째 bullet을 기본 패러그래프에 설정
                                tf.text = cleaned[0]
                                p0 = tf.paragraphs[0]
                                p0.level = 0
                                p0.font.name = "Calibri"
                                p0.font.size = Pt(32)
                                p0.font.color.rgb = body_color
                                # 나머지 bullet은 add_paragraph 로 추가
                                for b in cleaned[1:]:
                                    p = tf.add_paragraph()
                                    p.text = b
                                    p.level = 0
                                    p.font.name = "Calibri"
                                    p.font.size = Pt(32)
                                    p.font.color.rgb = body_color
                            

                            # 발표자 노트에 나레이션 추가
                            notes_slide = sld.notes_slide
                            notes_text_frame = notes_slide.notes_text_frame
                            notes_text_frame.text = f"발표자 나레이션: {slide.narration}"

                            # ─── DALL·E로 이미지 생성 및 삽입 ───────────────────────────
                            with st.spinner(f"이미지 생성 중: {slide.title}"):
                                # 테마에 따라 프롬프트 조정
                                theme_style = "dark, sleek" if theme_option == "다크 모드" else "light, clean"
                                prompt_text = (
                                    ## bullet 을 LLM 에게 줘서 핵심 요약을 뽑은 다음에 넣기기
                                    f"A minimalist, symbolic, and simplified vector-style graphic illustration representing {'; '.join(slide.title)}. Solid colors, clean lines, no text, suitable for professional presentation slides."
                                )
                                
                                img_resp = client.images.generate(
                                    model="dall-e-3",
                                    prompt=prompt_text,
                                    n=1,
                                    size="1024x1792", #1024X1792 오른쪽에 높이의 1/4(실험험) 만 차지하게 squeeze 해서 넣기
                                )
                                img_url = img_resp.data[0].url
                                img_data = requests.get(img_url).content
                                img_stream = io.BytesIO(img_data)

                                # 이미지를 더 오른쪽으로 치우기 랜덤하게게
                                right_width = Inches(8)  # 이미지 너비 조정
                                
                                # 오른쪽으로 더 밀기 위해 위치 조정 
                                left_position = prs.slide_width - right_width + Inches(4)  # 오른쪽으로 더 밀기
                                
                                # 이미지 높이 및 위치 조정
                                image_height = prs.slide_height - Inches(0)  # 여백 늘림
                                top_position = Inches(0)  # 상단 여백
                                
                                sld.shapes.add_picture(
                                    img_stream,
                                    left=left_position,
                                    top=top_position,
                                    width=right_width,
                                    height=image_height,
                                )

                        # PowerPoint 저장 - 모든 형식 변환에 사용할 버퍼
                        buf = io.BytesIO()
                        prs.save(buf)
                        buf.seek(0)
                        
                        # 생성된 PPT를 session_state에 저장
                        st.session_state.ppt_buffer = buf.getvalue()
                        
                        # 개별 다운로드 버튼 구현
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.download_button(
                                label="PowerPoint 다운로드 (.pptx)",
                                data=st.session_state.ppt_buffer,
                                file_name="presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            )
                            
                        with col2:
                            # PDF 다운로드 버튼
                            if st.button("PDF로 변환"):
                                with st.spinner("PDF로 변환 중..."):
                                    try:
                                        # 임시 디렉토리 생성
                                        import os
                                        import tempfile
                                        
                                        # 임시 디렉토리 생성
                                        temp_dir = tempfile.mkdtemp()
                                        
                                        # 임시 파일 경로 설정
                                        temp_pptx = os.path.join(temp_dir, "presentation.pptx")
                                        temp_pdf = os.path.join(temp_dir, "presentation.pdf")
                                        
                                        # PowerPoint 파일 저장
                                        with open(temp_pptx, "wb") as f:
                                            f.write(st.session_state.ppt_buffer)
                                        
                                        # pptxtopdf 모듈 임포트
                                        from pptxtopdf import convert
                                        
                                        # 변환 실행
                                        st.info("PowerPoint를 PDF로 변환하는 중...")
                                        convert(temp_pptx, temp_pdf)
                                        
                                        # 변환 성공 여부 확인
                                        if os.path.exists(temp_pdf):
                                            st.success("PDF 변환 완료!")
                                            
                                            # PDF 파일 읽기
                                            with open(temp_pdf, "rb") as f:
                                                pdf_data = f.read()
                                            
                                            # 세션 상태에 저장
                                            st.session_state.pdf_data = pdf_data
                                            st.session_state.pdf_ready = True
                                            
                                            # PDF 다운로드 버튼 표시
                                            st.download_button(
                                                label="PDF 다운로드",
                                                data=pdf_data,
                                                file_name="presentation.pdf",
                                                mime="application/pdf",
                                            )
                                            
                                            # 임시 파일 삭제
                                            import shutil
                                            shutil.rmtree(temp_dir)
                                        else:
                                            st.error("PDF 변환 실패 - 출력 파일을 찾을 수 없습니다")
                                    
                                    except ImportError:
                                        st.error("pptxtopdf 라이브러리를 찾을 수 없습니다")
                                        st.info("설치하려면: pip install pptxtopdf")
                                    except Exception as e:
                                        st.error(f"PDF 변환 중 오류 발생: {str(e)}")
                            
                            # 이미 변환된 PDF가 있으면 다운로드 버튼 표시
                            elif 'pdf_ready' in st.session_state and st.session_state.pdf_ready:
                                st.download_button(
                                    label="PDF 다운로드",
                                    data=st.session_state.pdf_data,
                                    file_name="presentation.pdf",
                                    mime="application/pdf",
                                )
                    except Exception as e:
                        st.error(f"PPT 생성 실패: {e}")

                else:
                    st.info("선택한 페이지에서 텍스트를 찾지 못했습니다.")
        except Exception as e:
            st.error(f"PDF를 읽는 중 오류가 발생했습니다: {e}")
    else:
        st.info("PDF 파일을 업로드해주세요.")

elif selected == "From Slide": # 강의 영상 생성
    # -*- coding: utf-8 -*-
    import os
    import sys
    import locale
    import uuid
    import tempfile
    import subprocess
    import io
    import base64

    import streamlit as st
    from pptx import Presentation
    import fitz                             # PyMuPDF for PDF text extraction
    import openai
    from openai import OpenAI
    from pdf2image import convert_from_path, convert_from_bytes
    from PIL import Image
    from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips

    import torch
    import torchaudio
    from zonos.model import Zonos, DEFAULT_BACKBONE_CLS as ZonosBackbone
    from zonos.conditioning import make_cond_dict, supported_language_codes
    from zonos.utils import DEFAULT_DEVICE as device

    from dotenv import load_dotenv

    # ─────────────── 초기 설정 ───────────────
    load_dotenv()  # .env 파일에 OPENAI_API_KEY=sk-... 설정

    # UTF-8 강제 인코딩
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
        sys.stdin .reconfigure(encoding="utf-8")
    os.environ["PYTHONUTF8"]       = "1"
    os.environ["PYTHONIOENCODING"] = "utf-8"
    locale.setlocale(locale.LC_ALL, 'C.UTF-8')

    # OpenAI 키 설정
    openai.api_key = os.getenv("OPENAI_API_KEY")
    if not openai.api_key or not openai.api_key.isascii():
        raise ValueError("유효한 OPENAI_API_KEY가 필요합니다.")

    # GPT-4o 클라이언트 생성
    client = OpenAI(api_key=openai.api_key)

    # Zonos TTS 기본 설정
    VOICE_NAME = "onyx"
    TTS_MODEL  = "tts-1-hd"

    # 스트림릿 페이지 설정
    #st.set_page_config(page_title="문서→자동 멀티모달 강의 비디오", layout="wide")
    st.title("📑 문서(.pptx/.pdf) → 🎬 멀티모달 자동 강의 비디오 생성기")

    # ─────────────── 유틸 함수 ───────────────

    def extract_from_pptx(data: bytes):
        prs = Presentation(io.BytesIO(data))
        slides = []
        for slide in prs.slides:
            # 텍스트 추출
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            text = "\n".join(text_parts).strip()
            # 이미지 추출
            imgs = []
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    # shape_type==13은 Picture
                    try:
                        blob = shape.image.blob
                        imgs.append(blob)
                    except Exception:
                        continue
            slides.append({"text": text, "images": imgs})
        return slides

    def extract_from_pdf(data: bytes):
        doc = fitz.open(stream=data, filetype="pdf")
        slides = []
        for page in doc:
            text = page.get_text("text", sort=True).strip()
            # 페이지 전체를 하나의 이미지로 변환
            pix = page.get_pixmap()
            img_bytes = pix.tobytes()
            slides.append({"text": text, "images": [img_bytes]})
        doc.close()
        return slides

    def convert_slides_to_images_auto(input_fp: str, out_dir: str, dpi: int = 200) -> list[str]:
        os.makedirs(out_dir, exist_ok=True)
        ext = os.path.splitext(input_fp)[1].lower()
        # PPTX → PDF
        if ext == ".pptx":
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir",
                os.path.dirname(input_fp), input_fp],
                check=True
            )
            pdf_path = os.path.splitext(input_fp)[0] + ".pdf"
        else:
            pdf_path = input_fp

        # PDF → PNG
        images = convert_from_path(pdf_path, dpi=dpi, fmt="png")
        saved = []
        for i, img in enumerate(images, start=1):
            fn = os.path.join(out_dir, f"slide_{i:03}.png")
            img.save(fn, "PNG")
            saved.append(fn)

        if ext == ".pptx" and os.path.exists(pdf_path):
            os.remove(pdf_path)
        return saved


    def call_gpt4o_with_history(text: str, image_bytes: bytes, scripts_so_far: list[str]) -> str:
        """
        GPT-4o에 텍스트와 이미지, 그리고 이전 슬라이드 스크립트들을 함께 전달하여
        한국어 내레이션 스크립트를 생성합니다.
        """
        system_msg = {
            "role": "system",
            "content": (
                "당신은 프레젠테이션 슬라이드를 간결하고 자연스럽게 설명하는 한국어 내레이션 스크립트 작성자입니다. "
                "끝맺음 인사말(예: ‘감사합니다’)이나 불필요한 여담 없이, 핵심 내용을 부드럽게 전달해주세요. "
                "이전 슬라이드에서 설명된 내용을 참고하여 자연스럽게 이어지도록 작성하세요."
            )
        }

        # 이전 슬라이드 스크립트들을 최대 2개까지만 포함
        history_content = []
        for prev_script in scripts_so_far[-2:]:
            history_content.append({"type": "text", "text": f"이전 슬라이드 스크립트: {prev_script}"})

        # 현재 슬라이드 텍스트
        content_list = history_content.copy()
        content_list.append({"type": "text", "text": text or ""})

        # 이미지가 있을 때만 image_url 블록 추가
        if image_bytes:
            b64 = base64.b64encode(image_bytes).decode("utf-8")
            content_list.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}"}
            })

        # 마지막으로 “스크립트 작성 요청” 텍스트
        content_list.append({
            "type": "text",
            "text": "위 슬라이드를 이전 내용과 자연스럽게 이어가며 한국어로 설명하는 간결한 스크립트를 작성해주세요."
        })

        user_msg = {"role": "user", "content": content_list}

        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[system_msg, user_msg]
        )

        script = resp.choices[0].message.content.strip()
        print(f"[GENERATED SCRIPT] Slide:\n{script}\n{'-'*40}")
        return script



    def save_uploaded_audio(uploaded_file) -> str | None:
        """
        업로드된 파일을 임시로 저장하고 파일 경로를 반환합니다.
        """
        if uploaded_file is None:
            return None
        suffix = os.path.splitext(uploaded_file.name)[1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uploaded_file.read())
            return tmp.name

    @st.cache_resource(show_spinner=False)
    def load_zonos_model(model_name: str) -> Zonos:
        model = Zonos.from_pretrained(model_name, device=device)
        model.requires_grad_(False).eval()
        return model

    def build_video(imgs: list[str], auds: list[str], out_mp4: str):
        clips = []
        for img_path, aud_path in zip(imgs, auds):
            ac = AudioFileClip(aud_path)
            ic = ImageClip(img_path).set_duration(ac.duration)
            clips.append(ic.set_audio(ac))
        final = concatenate_videoclips(clips, method="chain")
        final.write_videofile(out_mp4, codec="libx264", audio_codec="aac", fps=24)

    # ─────────────── 사이드바: 옵션 설정 ───────────────
    st.sidebar.header("옵션 설정")

    # 1) Voice Source 선택
    voice_source = st.sidebar.radio(
        "Voice Source 선택",
        options=["OpenAI TTS", "Voice Cloning (Zonos)"]
    )

    # 2) Zonos 설정 (Voice Cloning 선택 시)
    if voice_source == "Voice Cloning (Zonos)":
        # 지원 모델 확인
        supported_models = []
        if "transformer" in ZonosBackbone.supported_architectures:
            supported_models.append("Zyphra/Zonos-v0.1-transformer")
        if "hybrid" in ZonosBackbone.supported_architectures:
            supported_models.append("Zyphra/Zonos-v0.1-hybrid")
        if not supported_models:
            st.sidebar.error("지원되는 Zonos 모델이 없습니다.")
            st.stop()
        model_choice = st.sidebar.selectbox("Zonos 모델 선택", supported_models)

        # 화자 오디오 업로드
        ref_audio = st.sidebar.file_uploader(
            "화자 오디오 업로드 (wav)", type=["wav"], help="보이스 클로닝용 화자 음성"
        )

        # Conditioning Parameters
        st.sidebar.markdown("### Conditioning Parameters")
        dnsmos = st.sidebar.slider("DNSMOS Overall", min_value=1.0, max_value=5.0, value=4.0, step=0.1)
        fmax = st.sidebar.slider("Fmax (Hz)", min_value=0, max_value=24000, value=24000, step=1)
        vqscore = st.sidebar.slider("VQ Score", min_value=0.5, max_value=0.8, value=0.78, step=0.01)
        pitch_std = st.sidebar.slider("Pitch Std", min_value=0.0, max_value=300.0, value=45.0, step=1.0)
        speaking_rate = st.sidebar.slider("Speaking Rate", min_value=5.0, max_value=30.0, value=15.0, step=0.5)

        # Generation Parameters
        st.sidebar.markdown("### Generation Parameters")
        cfg_scale = st.sidebar.slider("CFG Scale", min_value=1.0, max_value=5.0, value=2.0, step=0.1)
        seed = st.sidebar.number_input("Seed", value=420, step=1, format="%d")
        randomize_seed = st.sidebar.checkbox("랜덤 시드 사용", value=True)

        # Sampling Parameters
        with st.sidebar.expander("Sampling Parameters"):
            linear = st.slider("Linear (0: 비활성)", min_value=-2.0, max_value=2.0, value=0.5, step=0.01)
            confidence = st.slider("Confidence", min_value=-2.0, max_value=2.0, value=0.40, step=0.01)
            quadratic = st.slider("Quadratic", min_value=-2.0, max_value=2.0, value=0.00, step=0.01)
            top_p = st.slider("Top P", min_value=0.0, max_value=1.0, value=0.0, step=0.01)
            top_k = st.slider("Top K", min_value=0, max_value=1024, value=0, step=1)
            min_p = st.slider("Min P", min_value=0.0, max_value=1.0, value=0.0, step=0.01)

        # Unconditional Keys
        with st.sidebar.expander("Advanced (Unconditional Keys)"):
            unconditional_keys = st.multiselect(
                "Unconditional Keys",
                options=[
                    "speaker", "emotion", "vqscore_8", "fmax",
                    "pitch_std", "speaking_rate", "dnsmos_ovrl", "speaker_noised"
                ],
                default=["emotion"],
                help="체크된 키는 조건 없이 자동으로 채워집니다."
            )

        # Emotion Sliders
        with st.sidebar.expander("Emotion Sliders"):
            e1 = st.slider("Happiness", min_value=0.0, max_value=1.0, value=1.0, step=0.05)
            e2 = st.slider("Sadness", min_value=0.0, max_value=1.0, value=0.05, step=0.05)
            e3 = st.slider("Disgust", min_value=0.0, max_value=1.0, value=0.05, step=0.05)
            e4 = st.slider("Fear", min_value=0.0, max_value=1.0, value=0.05, step=0.05)
            e5 = st.slider("Surprise", min_value=0.0, max_value=1.0, value=0.05, step=0.05)
            e6 = st.slider("Anger", min_value=0.0, max_value=1.0, value=0.05, step=0.05)
            e7 = st.slider("Other", min_value=0.0, max_value=1.0, value=0.10, step=0.05)
            e8 = st.slider("Neutral", min_value=0.0, max_value=1.0, value=0.20, step=0.05)
    else:
        model_choice = None
        ref_audio = None
        dnsmos = fmax = vqscore = pitch_std = speaking_rate = None
        cfg_scale = seed = None
        randomize_seed = False
        linear = confidence = quadratic = top_p = top_k = min_p = None
        unconditional_keys = ["emotion"]
        e1 = e2 = e3 = e4 = e5 = e6 = e7 = e8 = None

    # ─────────────── 초기 세션 상태 ───────────────
    if "video_ready" not in st.session_state:
        st.session_state.video_ready = False
        st.session_state.video_path  = ""

    # 파일 업로드
    uploaded = st.file_uploader("▶ PPTX 또는 PDF 파일 업로드", type=["pptx", "pdf"])
    if not uploaded:
        st.info("먼저 파일을 업로드하세요.")
        st.stop()

    data = uploaded.read()
    # 멀티모달 추출: 텍스트 + 이미지
    # 1) 업로드 파일 임시 저장
    temp_fp = os.path.join(tempfile.gettempdir(), uploaded.name)
    with open(temp_fp, "wb") as f:
        f.write(data)

    # 2) 슬라이드 전체 → PDF→PNG 변환
    img_dir = os.path.join(tempfile.gettempdir(), f"slides_{uuid.uuid4()}")
    slide_images = convert_slides_to_images_auto(temp_fp, img_dir)
    st.success(f"{len(slide_images)}개 슬라이드를 전체 이미지로 변환했습니다.")

    # 3) 텍스트 추출용 slides_data 준비
    if uploaded.name.lower().endswith(".pptx"):
        slides_data = extract_from_pptx(data)
    else:
        slides_data = extract_from_pdf(data)


    # 변환 시작 버튼
    if st.button("🔄 변환 시작", key="convert"):
        try:
            # 1) 슬라이드별 스크립트 생성 (GPT-4o 멀티모달)
            total_slides = len(slides_data)
            st.info(f"슬라이드별 스크립트 생성 중 (0/{total_slides})…")
            scripts = []
            progress_bar = st.progress(0)
            for i, slide in enumerate(slides_data, start=1):
                #st.info(f"슬라이드별 스크립트 생성 중 ({i}/{total_slides})…")
                text = slide["text"]
                img_bytes = slide["images"][0] if slide["images"] else b""
                script = call_gpt4o_with_history(text, img_bytes, scripts)
                scripts.append(script)
                progress_bar.progress(i / total_slides)
            st.write("✅ 스크립트 생성 완료")

            # 2) 오디오 합성
            st.info("2/6 음성 합성 중...")
            audio_files = []

            if voice_source == "OpenAI TTS":
                # OpenAI TTS 호출
                for i, sc in enumerate(scripts, start=1):
                    mp3_fp = os.path.join(tempfile.gettempdir(), f"slide_{i}.mp3")
                    result = openai.audio.speech.create(
                        model=TTS_MODEL,
                        voice=VOICE_NAME,
                        input=sc,
                        response_format="mp3"
                    )
                    with open(mp3_fp, "wb") as f:
                        f.write(result.content)
                    audio_files.append(mp3_fp)
                    st.progress(i / len(scripts))
            else:
                # Zonos Voice Cloning
                model_zonos = load_zonos_model(model_choice)

                if ref_audio is None:
                    st.error("Voice Cloning을 사용하려면 화자 오디오를 반드시 업로드하세요.")
                    st.stop()

                # 화자 오디오 로드 및 임베딩 생성
                ref_path = save_uploaded_audio(ref_audio)
                wav_ref, sr_ref = torchaudio.load(ref_path)
                speaker_emb = model_zonos.make_speaker_embedding(wav_ref, sr_ref)
                speaker_emb = speaker_emb.to(device, dtype=torch.bfloat16)

                emotion_tensor = torch.tensor([e1, e2, e3, e4, e5, e6, e7, e8], device=device)
                vq_tensor = torch.tensor([vqscore] * 8, device=device).unsqueeze(0)

                for i, sc in enumerate(scripts, start=1):
                    # cond_dict 구성 (한국어 스크립트)
                    cond_dict = make_cond_dict(
                        text=sc,
                        language="ko",
                        speaker=speaker_emb,
                        emotion=emotion_tensor,
                        vqscore_8=vq_tensor,
                        fmax=float(fmax),
                        pitch_std=float(pitch_std),
                        speaking_rate=float(speaking_rate),
                        dnsmos_ovrl=float(dnsmos),
                        speaker_noised=False,
                        device=device,
                        unconditional_keys=unconditional_keys,
                    )
                    conditioning = model_zonos.prepare_conditioning(cond_dict)

                    sampling_params = dict(
                        top_p=float(top_p),
                        top_k=int(top_k),
                        min_p=float(min_p),
                        linear=float(linear),
                        conf=float(confidence),
                        quad=float(quadratic),
                    )
                    if randomize_seed:
                        seed = torch.randint(0, 2**32 - 1, (1,)).item()
                    torch.manual_seed(int(seed))

                    # 고정된 충분한 토큰 수 (예: 4096)
                    max_new_tokens = 4096

                    with st.spinner(f"Zonos 생성 중… 슬라이드 {i}/{len(scripts)}"):
                        codes = model_zonos.generate(
                            prefix_conditioning=conditioning,
                            audio_prefix_codes=None,
                            max_new_tokens=max_new_tokens,
                            cfg_scale=float(cfg_scale),
                            batch_size=1,
                            sampling_params=sampling_params,
                        )

                    wav_out = model_zonos.autoencoder.decode(codes).cpu().detach()
                    # Tensor 차원 조정: 2D (channels, time)
                    if wav_out.dim() == 3:
                        wav_out = wav_out[0]
                    elif wav_out.dim() == 1:
                        wav_out = wav_out.unsqueeze(0)
                    if wav_out.dim() == 2 and wav_out.size(0) > 1:
                        wav_out = wav_out[0:1, :]

                    sr_out = model_zonos.autoencoder.sampling_rate
                    tmp_wav = os.path.join(tempfile.gettempdir(), f"slide_{i}.wav")
                    torchaudio.save(tmp_wav, wav_out, sr_out)

                    mp3_fp = os.path.join(tempfile.gettempdir(), f"slide_{i}.mp3")
                    subprocess.run(
                        ["ffmpeg", "-y", "-i", tmp_wav, "-codec:a", "libmp3lame", mp3_fp],
                        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True
                    )
                    audio_files.append(mp3_fp)
                    st.progress(i / len(scripts))

                os.remove(ref_path)

            st.write(f"✅ 음성 합성 완료 (Voice Source: {voice_source})")

            # 3) 이미지 파일 생성 (이미지 폴더에 저장)
            st.info("3/6 이미지 생성 중...")
            img_dir = os.path.join(tempfile.gettempdir(), f"imgs_{uuid.uuid4()}")
            os.makedirs(img_dir, exist_ok=True)
            saved_images = []
            for i, slide in enumerate(slides_data, start=1):
                # PPTX: 이미 분리된 이미지 블롭을 저장하거나, PDF: 페이지 이미지를 저장
                if slide["images"]:
                    img_bytes = slide["images"][0]
                    img = Image.open(io.BytesIO(img_bytes))
                    fn = os.path.join(img_dir, f"slide_{i}.png")
                    img.save(fn, "PNG")
                    saved_images.append(fn)
                else:
                    # 텍스트 전용 슬라이드를 위해 빈 배경 이미지 생성
                    blank = Image.new("RGB", (1280, 720), color=(255, 255, 255))
                    fn = os.path.join(img_dir, f"slide_{i}.png")
                    blank.save(fn, "PNG")
                    saved_images.append(fn)
                st.progress(i / len(slides_data))
            st.write(f"✅ {len(saved_images)}개 이미지 생성")

            # 4) 비디오 합성
            st.info("4/6 비디오 합성 중...")
            out_mp4 = os.path.join(tempfile.gettempdir(), f"lecture_{uuid.uuid4()}.mp4")
            build_video(slide_images, audio_files, out_mp4)
            st.write("✅ 비디오 합성 완료")

            # 5) 결과 저장 및 UI 업데이트
            st.session_state.video_ready = True
            st.session_state.video_path  = out_mp4

        except Exception as e:
            st.error(f"❌ 오류 발생: {e}")

    # 비디오가 생성된 경우 재생 및 다운로드 버튼
    if st.session_state.video_ready:
        st.success("🎉 비디오 생성 완료!")
        st.video(st.session_state.video_path)
        with open(st.session_state.video_path, "rb") as vf:
            st.download_button(
                label="📥 비디오 다운로드",
                data=vf,
                file_name="lecture.mp4",
                mime="video/mp4"
            )
    else:
        st.info("비디오가 아직 생성되지 않았습니다. 위의 버튼을 눌러 변환을 시작하세요.")
